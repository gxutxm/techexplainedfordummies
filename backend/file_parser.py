"""
file_parser.py — Extract plain text from uploaded files
=========================================================
Supported formats: PDF, DOCX, PPTX, TXT

Called by the /session/start-from-file endpoint.
Returns a plain string that gets passed straight into the interviewer agent
as source_text — exactly the same as if the user had typed it in.
"""

import io
from fastapi import UploadFile, HTTPException


SUPPORTED_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "text/plain",
}

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


async def extract_text(file: UploadFile) -> str:
    """
    Read an uploaded file and return its text content as a plain string.
    Raises HTTPException if the file type is unsupported or parsing fails.
    """
    filename = file.filename or ""
    ext = _get_extension(filename)

    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{ext}'. Allowed: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    contents = await file.read()

    try:
        if ext == ".pdf":
            return _parse_pdf(contents)
        elif ext == ".docx":
            return _parse_docx(contents)
        elif ext == ".pptx":
            return _parse_pptx(contents)
        elif ext == ".txt":
            return contents.decode("utf-8", errors="ignore").strip()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=422,
            detail=f"Could not parse file '{filename}': {str(e)}"
        )


def _get_extension(filename: str) -> str:
    import os
    return os.path.splitext(filename.lower())[1]


def _parse_pdf(contents: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(contents))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    if not text:
        raise HTTPException(status_code=422, detail="PDF appears to be scanned or image-only — no text could be extracted.")
    return text


def _parse_docx(contents: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(contents))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_pptx(contents: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(contents))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
    return "\n\n".join(slides_text)
